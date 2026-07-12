"""Loader: convierte cualquier formato soportado en texto plano.

Formatos: PDF, DOCX, PPTX, CSV, XLSX, JSON, TXT, HTML, Markdown.
"""
import csv
import json
from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


class UnsupportedFileTypeError(Exception):
    pass


def load_text(file_path: Path, extension: str) -> str:
    ext = extension.lower().lstrip(".")
    handler = _HANDLERS.get(ext)
    if handler is None:
        raise UnsupportedFileTypeError(f"Formato no soportado: .{ext}")
    return handler(file_path).strip()


def _load_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _load_markdown(path: Path) -> str:
    # Se indexa el markdown como texto; no requiere renderizado a HTML.
    return path.read_text(encoding="utf-8", errors="ignore")


def _load_html(path: Path) -> str:
    html = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text(separator="\n")


def _load_json(path: Path) -> str:
    data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    return json.dumps(data, ensure_ascii=False, indent=2)


def _load_csv(path: Path) -> str:
    lines = []
    with path.open(newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            lines.append(" | ".join(row))
    return "\n".join(lines)


def _load_xlsx(path: Path) -> str:
    wb = load_workbook(path, data_only=True, read_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"### Hoja: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def _load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _load_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _load_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"### Diapositiva {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


_HANDLERS = {
    "txt": _load_txt,
    "md": _load_markdown,
    "markdown": _load_markdown,
    "html": _load_html,
    "htm": _load_html,
    "json": _load_json,
    "csv": _load_csv,
    "xlsx": _load_xlsx,
    "pdf": _load_pdf,
    "docx": _load_docx,
    "pptx": _load_pptx,
}
